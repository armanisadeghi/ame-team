import json
from pptx import Presentation
from pptx.util import Inches
from aidream.settings import BASE_DIR


def json_to_presentation(json_data):
    prs = Presentation()

    for item in json_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = item['title']
        textbox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(6), Inches(4))
        text_frame = textbox.text_frame

        for point in item['points']:
            p = text_frame.add_paragraph()
            p.text = point

    prs.save('output.pptx')


if __name__ == '__main__':
    json_content = [
        {
            "title": "Introduction to Change and Disruption",
            "points": [
                "Calm Before the Storm: Often overlooked yet critical.",
                "Historical periods of calm preceding major technological shifts.",
                "Understanding change and disruption in the context of business and technology."
            ]
        },
        {
            "title": "Historical Technological Changes",
            "points": [
                "Examples: Internet, mobile phones, and social media platforms.",
                "Failure of established companies: Nokia, BlackBerry, AOL.",
                "Importance of recognizing impending market shifts."
            ]
        },
        {
            "title": "AI Impact on the Legal Industry",
            "points": [
                "Current period of calm before AI disruption in legal practice.",
                "Comparison to past technological disruptions in other industries.",
                "Urgency for legal professionals to prepare for AI-driven changes."
            ]
        },
        {
            "title": "Advancements in AI for Legal Practice",
            "points": [
                "Potential to develop an AI capable of passing the bar exam quickly.",
                "AI's future roles: Outperforming humans in certain legal tasks.",
                "Importance of recognizing AI's growing capabilities in the legal field."
            ]
        },
        {
            "title": "Adapting to New Technologies in Business",
            "points": [
                "Personal technology choices vs. business technology decisions.",
                "Risks of obsolescence by not adopting new technologies.",
                "Staying competitive by embracing AI advancements."
            ]
        },
        {
            "title": "Revolutionizing Professional Practices with AI",
            "points": [
                "Introduction of a new AI application for professional use.",
                "No-code functionality for performing complex tasks.",
                "Empowering professionals to work more efficiently."
            ]
        },
        {
            "title": "Personal Journey with AI",
            "points": [
                "From no programming experience to sophisticated AI use.",
                "Synergies between human expertise and AI capabilities.",
                "Potential for transformative outcomes in various professions."
            ]
        },
        {
            "title": "Impact of AI on the Workforce",
            "points": [
                "Differentiating facts, opinions, and well-researched information.",
                "Significant job impact worldwide due to AI.",
                "Research needed to confirm specific industry numbers."
            ]
        },
        {
            "title": "Job Consolidation Due to AI",
            "points": [
                "Exploring how AI will change job execution, not eliminate jobs.",
                "Analogy of aircraft capacity to illustrate job consolidation.",
                "Adaptation through AI for role retention and efficiency."
            ]
        },
        {
            "title": "Future of AI in Industries",
            "points": [
                "AI enabling more productivity with less effort.",
                "Competitive landscape intensifying due to shared AI tools.",
                "Importance of staying current with AI advancements."
            ]
        }
    ]

    json_to_presentation(json_content)