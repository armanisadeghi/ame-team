import json
import os
import uuid

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from ame_team.settings.base import BASE_DIR
from common import get_sample_data, pretty_print, print_link

output_temp_dir = "temp/app_outputs/slides"
output_dir = os.path.join(BASE_DIR, output_temp_dir)
short_uuid_part = lambda: str(uuid.uuid4())[:8]


def setup_title(slide, title, params):
    title_shape = slide.shapes.title
    title_text_frame = title_shape.text_frame
    title_text_frame.text = title
    for paragraph in title_text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = params.get('title_bold', True)
            run.font.size = Pt(params.get('title_font_size', 32))


def add_bullet_points(slide, points, params):
    content_shape = slide.placeholders[1]  # Assumes placeholder 1 is the content placeholder
    text_frame = content_shape.text_frame
    text_frame.clear()  # Clear existing content in the placeholder
    for point in points:
        p = text_frame.add_paragraph()
        p.text = point
        p.level = 0  # Top level bullet points
        p.font.size = Pt(params.get('bullet_font_size', 24))


def set_background_color(slide, params):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill_color = params.get('background_color', (173, 216, 230))  # Default light blue
    fill.fore_color.rgb = RGBColor(*fill_color)


def json_to_presentation(json_data, file_name_root='output', output_dir='.', ppt_params={}):
    prs = Presentation()

    for item in json_data:
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Assumes layout 1 with title and content placeholder
        setup_title(slide, item['title'], ppt_params)
        add_bullet_points(slide, item['points'], ppt_params)
        set_background_color(slide, ppt_params)

    output_file_name = f"{file_name_root}_{short_uuid_part()}.pptx"
    save_path = os.path.join(output_dir, output_file_name)
    prs.save(save_path)
    print_link(save_path)
    return save_path


if __name__ == '__main__':
    sample_json = get_sample_data(app_name='automation_matrix', data_name='title_three_bullets_sample_json', sub_app='slide_source_samples')  # Get sample API response
    pretty_print(sample_json)

    ppt_params = {
        'title_font_size': 36,
        'bullet_font_size': 28,
        'background_color': (200, 220, 240)  # A different shade of light blue
    }

    save_path = json_to_presentation(sample_json, ppt_params=ppt_params)
    print("Presentation saved to:", save_path)
