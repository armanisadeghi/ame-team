from pptx import Presentation
import xml.etree.ElementTree as ET

# Load the presentation template
template_path = "template.pptx"
prs = Presentation(template_path)


# Example function to replace text placeholders in XML
def replace_placeholder(xml_file, placeholder, new_text):
    tree = ET.parse(xml_file)
    root = tree.getroot()

    for elem in root.iter():
        if elem.text and placeholder in elem.text:
            elem.text = elem.text.replace(placeholder, new_text)

    tree.write(xml_file)


# Example: Replace placeholders in the first slide
slide = prs.slides[0]
slide_xml_path = slide.partname[1:]  # Get the XML path of the slide

replace_placeholder(slide_xml_path, "{{Title}}", "New Title")
replace_placeholder(slide_xml_path, "{{Subtitle}}", "New Subtitle")

# Save the modified presentation
output_path = "output.pptx"
prs.save(output_path)
